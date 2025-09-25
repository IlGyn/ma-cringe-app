import streamlit as st
import numpy as np
import json
import io
import csv
import docx
import uuid
import torch
import requests
import socket
import aiohttp
import asyncio
import logging
import threading
from queue import Queue
from datetime import datetime
from PyPDF2 import PdfReader
from pptx import Presentation

# Добавляем fallback для PDF и других библиотек
try:
    import pdfplumber
    PDF_PLUMBER_AVAILABLE = True
except ImportError:
    PDF_PLUMBER_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
from qdrant_client.http import models as qmodels
from typing import List, Dict, Any, AsyncGenerator, Tuple, Union
from collections import OrderedDict
from urllib.parse import urlparse
from embedding_cache import EmbeddingCache
import os
import html

# ==================== CONFIGURATION ====================
class Config:
    QDRANT_HOST = os.getenv("QDRANT_HOST", "localhost")
    QDRANT_PORT = int(os.getenv("QDRANT_PORT", 6334))
    COLLECTION_NAME = "chat_memory"
    OLLAMA_URL = os.getenv("OLLAMA_URL", "http://localhost:11434")
    MAX_FILE_SIZE_MB = 10
    MAX_CACHE_SIZE = 100
    DEFAULT_EMBEDDING_MODEL = "all-MiniLM-L6-v2"
    DEFAULT_LLM_MODEL = "Phi4-mini"
    DEBUG_MODE = os.getenv("DEBUG", "false").lower() == "true"
    EMBEDDING_CACHE_SIZE = int(os.getenv("EMBEDDING_CACHE_SIZE", 500))
    ASYNC_EMBEDDING_ENABLED = os.getenv("ASYNC_EMBEDDING_ENABLED", "true").lower() == "true"
    
    # Новые настройки
    MAX_MESSAGE_LENGTH = 5000
    CHUNK_SIZE = 800
    CHUNK_OVERLAP = 100
    MAX_CONCURRENT_UPLOADS = 3
    STREAM_TIMEOUT = 120  # Увеличенный таймаут для стриминга

# ==================== LOGGING ====================
logging.basicConfig(
    filename="chat_app.log",
    level=logging.INFO if not Config.DEBUG_MODE else logging.DEBUG,
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s"
)

# ==================== CORE CLASSES ====================
class MessageManager:
    @staticmethod
    def load_messages(qdrant_manager, limit: int = 50) -> List[Dict[str, Any]]:
        try:
            hits = qdrant_manager.client.scroll(
                collection_name=qdrant_manager.collection_name, 
                limit=limit, 
                with_payload=True, 
                with_vectors=False
            )
            return [h.payload for h in hits[0] if h.payload is not None] if hits and hits[0] else []
        except Exception as e:
            logging.error(f"Ошибка загрузки сообщений: {e}")
            return []
    
    @staticmethod
    def truncate_context_by_tokens(context: List[Dict], max_tokens: int = 2000) -> List[Dict]:
        """Обрезает контекст по приблизительному количеству токенов"""
        total_chars = 0
        truncated = []
        
        for msg in reversed(context):
            content = msg.get('content', '')
            chars = len(content)
            if total_chars + chars > max_tokens * 4:
                break
            truncated.append(msg)
            total_chars += chars
        
        return list(reversed(truncated))

    @staticmethod
    def is_suspicious(content: str) -> bool:
        """Проверяет, содержит ли сообщение подозрительные фразы"""
        suspicious_phrases = [
            'согласно предыдущим взаимодействиям',
            'уже общаемся часто', 
            'знакомы с моими возможностями',
            'новый взгляд может быть очень ценным',
            'согласно последнему взаимодействию',
            'я только был загружен',
            'не имею воспоминаний'
        ]
        content_lower = content.lower()
        return any(pattern in content_lower for pattern in suspicious_phrases)

class QdrantManager:
    def __init__(self, host: str, port: int, collection_name: str):
        self.client = QdrantClient(host=host, port=port, prefer_grpc=True)
        self.collection_name = collection_name

    def recreate_collection(self, vector_size: int):
        try:
            self.client.delete_collection(self.collection_name)
            logging.info(f"Коллекция {self.collection_name} удалена")
        except Exception as e:
            logging.warning(f"Не удалось удалить коллекцию: {e}")
        
        self.client.create_collection(
            collection_name=self.collection_name,
            vectors_config=qmodels.VectorParams(size=vector_size, distance=qmodels.Distance.COSINE)
        )
        logging.info(f"Коллекция {self.collection_name} создана")

    def init_collection(self, encoder):
        vector_size = encoder.get_sentence_embedding_dimension()
        try:
            collections = [c.name for c in self.client.get_collections().collections]
            if self.collection_name not in collections:
                self.recreate_collection(vector_size)
                return

            coll_info = self.client.get_collection(self.collection_name)
            coll_vectors = getattr(coll_info.config.params, 'vectors', None)

            recreate_needed = False
            if isinstance(coll_vectors, qmodels.VectorParams):
                if coll_vectors.size != vector_size:
                    recreate_needed = True
            elif isinstance(coll_vectors, dict):
                first_vec = next(iter(coll_vectors.values()), None)
                if first_vec and getattr(first_vec, 'size', None) != vector_size:
                    recreate_needed = True
            else:
                recreate_needed = True

            if recreate_needed:
                self.recreate_collection(vector_size)

        except Exception as e:
            logging.error(f"Ошибка инициализации коллекции: {e}")
            self.recreate_collection(vector_size)

    def save_message(self, role: str, content: str, timestamp: str, encoder_manager):
        # ФИЛЬТР: не сохраняем странные ответы
        if MessageManager.is_suspicious(content):
            logging.warning(f"Не сохраняем подозрительное сообщение: {content[:100]}...")
            return
        
        try:
            vector = encoder_manager.get_embedding(content)
            point = qmodels.PointStruct(
                id=str(uuid.uuid4()),
                vector=vector,
                payload={"role": role, "content": content, "time": timestamp}
            )
            self.client.upsert(collection_name=self.collection_name, points=[point])
            logging.info(f"Сообщение сохранено: {role}, {timestamp}")
        except Exception as e:
            logging.error(f"Ошибка сохранения в Qdrant: {e}")

    def search_context(self, prompt: str, encoder_manager, limit: int = 6) -> List[Dict[str, Any]]:
        # Проверка кэша
        cache_key = f"{prompt}_{limit}"
        if "context_cache" not in st.session_state:
            st.session_state.context_cache = OrderedDict()
        
        cached = st.session_state.context_cache.get(cache_key)
        if cached is not None:
            logging.info(f"Контекст для '{prompt}' загружен из кэша")
            return cached
        
        try:
            vector = encoder_manager.get_embedding(prompt)
  
            response = self.client.query_points(
                collection_name=self.collection_name, 
                query=vector, 
                limit=limit * 2,
                with_payload=True
            )
            
            all_results = [p.payload for p in getattr(response, 'points', [])] if getattr(response, 'points', None) is not None else []
            
            chat_messages = [msg for msg in all_results if msg.get('role') in ['user', 'assistant']]
            file_chunks = [msg for msg in all_results if msg.get('role') == 'file']
            
            recent_chat = chat_messages[-3:] if len(chat_messages) > 3 else chat_messages
            relevant_files = file_chunks[:3]
            
            result = recent_chat + relevant_files
            
            # ФИЛЬТРАЦИЯ: убираем подозрительные сообщения
            filtered_result = [msg for msg in result if not MessageManager.is_suspicious(msg.get('content', ''))]
            
            if len(filtered_result) != len(result):
                logging.info(f"Отфильтровано {len(result) - len(filtered_result)} сообщений из поиска")
            
            # Сохраняем в кэш
            st.session_state.context_cache[cache_key] = filtered_result
            if len(st.session_state.context_cache) > Config.MAX_CACHE_SIZE:
                st.session_state.context_cache.popitem(last=False)
            
            return filtered_result
        except Exception as e:
            logging.error(f"Ошибка поиска в Qdrant: {e}")
            return []

    def load_messages(self, limit: int = 50) -> List[Dict[str, Any]]:
        """Загрузка всех сообщений из коллекции"""
        try:
            hits = self.client.scroll(
                collection_name=self.collection_name, 
                limit=limit, 
                with_payload=True, 
                with_vectors=False
            )
            return [h.payload for h in hits[0] if h.payload is not None] if hits and hits[0] else []
        except Exception as e:
            logging.error(f"Ошибка загрузки сообщений: {e}")
            return []

    def clear_collection(self, encoder):
        try:
            self.client.delete_collection(self.collection_name)
            self.recreate_collection(encoder.get_sentence_embedding_dimension())
        except Exception as e:
            logging.error(f"Ошибка очистки коллекции: {e}")

@st.cache_resource(show_spinner=False)
def load_sentence_transformer(model_name: str = "all-MiniLM-L6-v2"):
    import torch
    device = "cuda" if torch.cuda.is_available() else "cpu"
    logging.info(f"Загрузка модели эмбеддингов: {model_name} на {device}")
    return SentenceTransformer(model_name, device=device)

class EncoderManager:
    def __init__(self):
        self.cache = EmbeddingCache(max_size=Config.EMBEDDING_CACHE_SIZE)
        self._encoder = None

    def load_encoder(self, model_name: str = "all-MiniLM-L6-v2"):
        return load_sentence_transformer(model_name)

    def get_encoder(self):
        if self._encoder is None:
            self._encoder = self.load_encoder(Config.DEFAULT_EMBEDDING_MODEL)
        return self._encoder

    def get_embedding(self, text: str, encoder=None) -> List[float]:
        if encoder is None:
            encoder = self.get_encoder()
            
        # Проверка кэша
        cached = self.cache.get(text)
        if cached is not None:
            logging.debug(f"Эмбеддинг из кэша: {text[:50]}...")
            return cached
        
        # Генерация эмбеддинга
        vec = encoder.encode(text, convert_to_numpy=False)
        vec_list: List[float] = []
        
        if hasattr(vec, 'tolist'):
            vec_list = [float(x) for x in vec.tolist()]
        elif isinstance(vec, (list, np.ndarray)):
            vec_list = [float(x) for x in vec]
        else:
            vec_list = [float(vec)]

        # Сохраняем в кэш
        self.cache.put(text, vec_list)
        return vec_list

    async def async_encode(self, texts: List[str], encoder=None, batch_size: int = 32) -> List[List[float]]:
        """Асинхронная генерация эмбеддингов"""
        if encoder is None:
            encoder = self.get_encoder()
            
        loop = asyncio.get_event_loop()
        
        # Проверяем кэш для каждого текста
        cached_vectors: List[tuple] = []
        texts_to_encode: List[str] = []
        text_indices: List[int] = []
        
        for i, text in enumerate(texts):
            cached = self.cache.get(text)
            if cached is not None:
                cached_vectors.append((i, cached))
                logging.debug(f"Эмбеддинг из кэша: {text[:50]}...")
            else:
                texts_to_encode.append(text)
                text_indices.append(i)
        
        result_vectors: List[Union[List[float], None]] = [None] * len(texts)
        
        if texts_to_encode:
            vectors = await loop.run_in_executor(
                None, 
                lambda: encoder.encode(texts_to_encode, batch_size=batch_size, convert_to_numpy=False)
            )
            
            if hasattr(vectors, 'tolist'):
                vectors = vectors.tolist()
            elif not isinstance(vectors, list):
                vectors = list(vectors)
            
            processed_vectors: List[List[float]] = []
            for vec in vectors:
                if hasattr(vec, 'tolist'):
                    vec = vec.tolist()
                if isinstance(vec, (list, np.ndarray)):
                    processed_vectors.append([float(x) for x in vec])
                else:
                    processed_vectors.append([float(vec)])
            
            for idx, vec in cached_vectors:
                result_vectors[idx] = vec
            
            for i, vec in enumerate(processed_vectors):
                text_idx = text_indices[i]
                result_vectors[text_idx] = vec
                self.cache.put(texts_to_encode[i], vec)
        else:
            for idx, vec in cached_vectors:
                result_vectors[idx] = vec
        
        final_result: List[List[float]] = []
        for vec in result_vectors:
            if vec is not None:
                final_result.append(vec)
        
        return final_result

class OllamaClient:
    def check_connection(self, host: str = "127.0.0.1", port: int = 11434, timeout: int = 2) -> bool:
        try:
            with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
                sock.settimeout(timeout)
                return sock.connect_ex((host, port)) == 0
        except Exception as e:
            logging.error(f"Ошибка проверки порта: {e}")
            return False

    def check_api(self, base_url: str) -> Tuple[bool, str]:
        if not self.check_connection():
            return False, "Порт 11434 закрыт"
        try:
            r = requests.get(f"{base_url}/api/version", timeout=2)
            if r.status_code == 200:
                return True, "Подключение установлено"
            else:
                return False, "Нет ответа от сервера"
        except requests.exceptions.RequestException as e:
            return False, str(e)

    async def ask_llm_async_stream(self, prompt: str, context: list, model: str, 
                                    base_url: str, max_tokens: int = 200,
                                    temperature: float = 0.3, top_p: float = 0.8, 
                                    top_k: int = 40) -> AsyncGenerator[str, None]:

        system_instruction = (
            "Ты — виртуальный помощник. Отвечай по существу, используя предоставленный контекст из файлов и истории чата.\n"
            "Если в контексте есть релевантная информация из загруженных файлов, используй её (например, 'Согласно загруженному файлу...').\n"
            "Если контекста нет или он не релевантен, отвечай на основе общего знания.\n"
            "НЕ добавляй фразы типа 'Согласно предыдущим взаимодействиям' или 'Как дела у вас?'.\n"
            "Отвечай напрямую на вопрос пользователя."
        )
        
        # Фильтрация контекста
        filtered_context = [msg for msg in context if not MessageManager.is_suspicious(msg.get('content', ''))]
        
        if len(filtered_context) != len(context):
            logging.info(f"Отфильтровано {len(context) - len(filtered_context)} сообщений для prompt")
            context = filtered_context
        
        full_prompt = "\n".join([
            f"{m['role'].capitalize()}: {m['content']}" 
            for m in context 
            if m.get('content') and m.get('content').strip()
        ] + [f"User: {prompt}", "Assistant:"])
        
        data = {
            "model": model,
            "prompt": f"System: {system_instruction}\n{full_prompt}",
            "stream": True,
            "options": {
                "num_predict": max_tokens,
                "temperature": temperature,
                "top_p": top_p,
                "top_k": top_k,
                "num_ctx": 8192
            }
        }
        
        async with aiohttp.ClientSession() as session:
            try:
                timeout = aiohttp.ClientTimeout(total=Config.STREAM_TIMEOUT, connect=30)
                async with session.post(f"{base_url}/api/generate", json=data, timeout=timeout) as r:
                    if r.status != 200:
                        error_text = await r.text()
                        yield f"⚠️ Ошибка Ollama ({r.status}): {error_text}"
                        return
                        
                    async for raw_line in r.content:
                        if not raw_line:
                            continue
                        line = raw_line.strip()
                        if not line:
                            continue
                        
                        try:
                            line_str = line.decode('utf-8')
                            if line_str.startswith('data:'):
                                json_str = line_str[5:].strip()
                            else:
                                json_str = line_str
                                
                            if json_str:
                                chunk = json.loads(json_str)
                                if 'error' in chunk:
                                    yield f"⚠️ Ошибка модели: {chunk['error']}"
                                    break
                                elif chunk.get("response"):
                                    yield chunk["response"]
                                elif chunk.get("done"):
                                    break
                        except json.JSONDecodeError as je:
                            logging.warning(f"Не JSON данные: {line_str[:100]}...")
                            continue
                        except Exception as e:
                            logging.error(f"Ошибка обработки chunk: {e}")
                            yield f"⚠️ Ошибка обработки: {e}"
                            break
            except asyncio.TimeoutError:
                yield "⚠️ Таймаут при запросе к модели"
            except Exception as e:
                error_msg = f"⚠️ Ошибка Ollama: {str(e)}"
                logging.error(f"Ошибка Ollama: {e}")
                yield error_msg

class FileProcessor:
    def __init__(self):
        pass

    def chunk_text(self, text: str, chunk_size: int = 800, overlap: int = 100) -> List[str]:
        if not text:
            return []
        chunks = []
        start = 0
        text_len = len(text)
        while start < text_len:
            end = min(start + chunk_size, text_len)
            chunks.append(text[start:end])
            start += chunk_size - overlap
            if start >= text_len:
                break
        return chunks

    def process_file(self, file, qdrant_manager: QdrantManager, encoder_manager, max_size_mb: int = 10) -> str:
        file_ext = file.name.lower().split('.')[-1] if '.' in file.name else ""
        allowed_extensions = {"txt", "pdf", "docx", "csv"}
        
        # Добавляем поддержку новых форматов только если библиотеки доступны
        if PPTX_AVAILABLE:
            allowed_extensions.add("pptx")
        if XLSX_AVAILABLE:
            allowed_extensions.add("xlsx")
        
        if not file_ext:
            return f"⚠️ Файл без расширения не поддерживается"
        
        if file_ext not in allowed_extensions:
            available_formats = ", ".join(sorted(allowed_extensions))
            return f"⚠️ Формат .{file_ext} не поддерживается. Доступные: {available_formats}"
        
        if file.size > max_size_mb * 1024 * 1024:
            return f"⚠️ Файл '{file.name}' слишком большой (макс. {max_size_mb} МБ)"

        try:
            file.seek(0)
            raw_bytes = file.read()
            file.seek(0)
        except Exception as e:
            return f"⚠️ Не удалось прочитать файл {file.name}: {e}"

        content = ""
        try:
            if file_ext == "txt":
                content = raw_bytes.decode("utf-8", errors="ignore")
            elif file_ext == "pdf":
                content = self._extract_pdf_content(raw_bytes)
            elif file_ext == "docx":
                content = self._extract_docx_content(raw_bytes)
            elif file_ext == "csv":
                content = self._extract_csv_content(raw_bytes)
            elif file_ext == "pptx":
                if PPTX_AVAILABLE:
                    content = self._extract_pptx_content(raw_bytes)
                else:
                    content = "⚠️ Обработка PPTX не поддерживается (отсутствует библиотека python-pptx)"
            elif file_ext == "xlsx":
                if XLSX_AVAILABLE:
                    content = self._extract_xlsx_content(raw_bytes)
                else:
                    content = "⚠️ Обработка XLSX не поддерживается (отсутствует библиотека openpyxl)"
            else:
                content = f"⚠️ Обработка формата .{file_ext} не поддерживается"
        except Exception as e:
            content = f"⚠️ Ошибка при обработке файла: {e}"
            logging.error(f"Ошибка обработки файла {file.name}: {e}")

        chunks = self.chunk_text(content)
        if chunks:
            self._save_chunks_as_messages(chunks, datetime.now().strftime("%H:%M:%S"), qdrant_manager, encoder_manager)
            if "stats" in st.session_state:
                st.session_state.stats["files_processed"] = st.session_state.stats.get("files_processed", 0) + 1
                st.session_state.stats["chunks_saved"] = st.session_state.stats.get("chunks_saved", 0) + len(chunks)
            return f"✅ Файл '{file.name}' загружен ({len(chunks)} чанков)"
        else:
            return f"⚠️ Файл '{file.name}' пустой или нераспознанный"

    def _extract_pdf_content(self, raw_bytes) -> str:
        # Пробуем PyPDF2
        try:
            reader = PdfReader(io.BytesIO(raw_bytes))
            total_pages = len(reader.pages)

            if total_pages == 0:
                return ""
            
            progress_bar = st.progress(0) if total_pages > 2 else None
            
            content_parts = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    content_parts.append(text + "\n")
                else:
                    # Fallback: пробуем pdfplumber если страница пустая
                    if PDF_PLUMBER_AVAILABLE:
                        try:
                            with pdfplumber.open(io.BytesIO(raw_bytes)) as pdf:
                                if i < len(pdf.pages):
                                    page_obj = pdf.pages[i]
                                    text_alt = page_obj.extract_text()
                                    if text_alt:
                                        content_parts.append(text_alt + "\n")
                        except Exception:
                            pass
                
                if progress_bar and i % max(1, total_pages // 20) == 0:
                    progress_bar.progress((i + 1) / total_pages)
            
            if progress_bar:
                progress_bar.empty()
            
            return "".join(content_parts)
        except Exception as e:
            logging.warning(f"PyPDF2 ошибка: {e}")
            # Fallback на pdfplumber если доступен
            if PDF_PLUMBER_AVAILABLE:
                try:
                    content_parts = []
                    with pdfplumber.open(io.BytesIO(raw_bytes)) as pdf:
                        progress_bar = st.progress(0) if len(pdf.pages) > 2 else None
                        for i, page in enumerate(pdf.pages):
                            text = page.extract_text()
                            if text:
                                content_parts.append(text + "\n")
                            if progress_bar and i % max(1, len(pdf.pages) // 20) == 0:
                                progress_bar.progress((i + 1) / len(pdf.pages))
                        if progress_bar:
                            progress_bar.empty()
                    return "".join(content_parts)
                except Exception as e2:
                    logging.error(f"pdfplumber ошибка: {e2}")
            return f"⚠️ Ошибка извлечения PDF: {e}"

    def _extract_docx_content(self, raw_bytes) -> str:
        doc = docx.Document(io.BytesIO(raw_bytes))
        content_parts = [para.text + "\n" for para in doc.paragraphs]
        return "".join(content_parts)

    def _extract_csv_content(self, raw_bytes) -> str:
        decoded = io.StringIO(raw_bytes.decode("utf-8", errors="ignore"))
        reader = csv.reader(decoded)
        return "\n".join([", ".join(row) for row in reader])

    def _extract_pptx_content(self, file_path: str) -> str:
        prs = Presentation(file_path)
        content_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            content_parts.append(f"Slide {slide_num}:\n")
            for shape in slide.shapes:
                # Извлекаем текст из text_frame, если он есть
                if getattr(shape, "has_text_frame", False):
                    text_frame = getattr(shape, "text_frame", None)
                    if text_frame:
                        for paragraph in getattr(text_frame, "paragraphs", []):
                            text = getattr(paragraph, "text", "").strip()
                            if text:
                                content_parts.append(text + "\n")
                else:
                    text = getattr(shape, "text", "").strip()
                    if text:
                        content_parts.append(text + "\n")

        return "\n".join(content_parts)


    def _extract_xlsx_content(self, raw_bytes) -> str:
        if not XLSX_AVAILABLE:
            return "⚠️ Библиотека openpyxl не установлена"
            
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True)
            content_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_parts.append(f"Sheet: {sheet_name}\n")
                for row in sheet.iter_rows(values_only=True):
                    content_parts.append(", ".join(str(cell) if cell is not None else "" for cell in row) + "\n")
            return "".join(content_parts)
        except Exception as e:
            return f"⚠️ Ошибка извлечения XLSX: {e}"

    async def _save_chunks_as_messages_async(self, chunks: list, timestamp: str, 
                                           qdrant_manager: QdrantManager, encoder_manager, 
                                           role: str = "file", progress_bar=None):
        try:
            if not chunks:
                return
                
            vectors = await encoder_manager.async_encode(chunks, batch_size=min(32, len(chunks)))
            
            points = []
            for chunk, vec in zip(chunks, vectors):
                vector_list = vec if isinstance(vec, list) else list(vec)
                point = qmodels.PointStruct(
                    id=str(uuid.uuid4()),
                    vector=vector_list,
                    payload={"role": role, "content": chunk, "time": timestamp}
                )
                points.append(point)
            
            batch_size = 64
            total_points = len(points)
            
            for i in range(0, total_points, batch_size):
                batch = points[i:i + batch_size]
                qdrant_manager.client.upsert(
                    collection_name=qdrant_manager.collection_name, 
                    points=batch
                )
                
                if progress_bar and total_points > batch_size:
                    progress_bar.progress(min(1.0, (i + batch_size) / total_points))
            
            logging.info(f"Сохранено {len(points)} чанков текста (асинхронно)")
            
        except Exception as e:
            logging.error(f"Ошибка асинхронного сохранения чанков: {e}")
            raise

    def _save_chunks_as_messages(self, chunks: list, timestamp: str, 
                               qdrant_manager: QdrantManager, encoder_manager, role: str = "file"):
        try:
            try:
                loop = asyncio.get_event_loop()
            except RuntimeError:
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)
            
            progress = st.progress(0) if len(chunks) > 32 else None
            
            # Используем asyncio.wait_for для таймаута
            try:
                future = asyncio.wait_for(
                    self._save_chunks_as_messages_async(chunks, timestamp, qdrant_manager, encoder_manager, role, progress),
                    timeout=Config.STREAM_TIMEOUT
                )
                loop.run_until_complete(future)
            except asyncio.TimeoutError:
                st.sidebar.warning("⚠️ Таймаут при сохранении чанков")
                logging.error("Таймаут при сохранении чанков")
            
            if progress:
                progress.empty()
                
        except Exception as e:
            logging.error(f"Ошибка сохранения чанков: {e}")
            st.sidebar.warning(f"Ошибка сохранения чанков: {e}")

class Validator:
    def validate_prompt(self, prompt: str) -> bool:
        if not prompt or not prompt.strip():
            st.error("Пожалуйста, введите непустое сообщение")
            logging.warning("Попытка отправки пустого сообщения")
            return False
        
        if len(prompt) > Config.MAX_MESSAGE_LENGTH:
            st.error(f"Сообщение слишком длинное (максимум {Config.MAX_MESSAGE_LENGTH} символов)")
            return False

        # Расширенная проверка безопасности
        dangerous_patterns = [
            "<?php", "<script", "javascript:", "onload=", 
            "base64_decode", "eval(", "system(", "exec(",
            "union select", "drop table", "insert into",
            "delete from", "update.*set", "create table"
        ]
        
        # HTML escape для дополнительной безопасности
        prompt_sanitized = html.escape(prompt.lower())
        
        if any(pattern in prompt_sanitized for pattern in dangerous_patterns):
            st.error("Сообщение содержит подозрительные паттерны")
            logging.warning(f"Обнаружены опасные паттерны в сообщении: {prompt[:100]}...")
            return False
            
        # Дополнительная проверка на SQL-инъекции
        sql_keywords = ['select', 'insert', 'update', 'delete', 'drop', 'create', 'alter']
        if any(keyword in prompt_sanitized for keyword in sql_keywords):
            # Проверяем на подозрительные комбинации
            suspicious_combos = ['\'', '"', '--', '/*', '*/', ';']
            if any(combo in prompt for combo in suspicious_combos):
                st.error("Обнаружена возможная SQL-инъекция")
                logging.warning(f"Возможная SQL-инъекция: {prompt[:100]}...")
                return False
            
        return True

    def validate_url(self, url: str) -> bool:
        if not url:
            return False
            
        try:
            if not url.startswith(('http://', 'https://')):
                return False
                
            result = urlparse(url)
            if not all([result.scheme, result.netloc]):
                return False
                
            return True
        except Exception:
            return False

    def validate_file_extension(self, filename: str) -> bool:
        allowed_extensions = {"txt", "pdf", "docx", "csv"}
        if PPTX_AVAILABLE:
            allowed_extensions.add("pptx")
        if XLSX_AVAILABLE:
            allowed_extensions.add("xlsx")
        file_ext = filename.lower().split('.')[-1] if '.' in filename else ""
        return file_ext in allowed_extensions

# ==================== GLOBAL INSTANCES ====================
config = Config()
qdrant_manager = QdrantManager(config.QDRANT_HOST, config.QDRANT_PORT, config.COLLECTION_NAME)
encoder_manager = EncoderManager()
ollama_client = OllamaClient()
file_processor = FileProcessor()
validator = Validator()
message_manager = MessageManager()

# ==================== SESSION STATE INITIALIZATION ====================
if "encoder" not in st.session_state:
    st.session_state.encoder = encoder_manager.load_encoder(config.DEFAULT_EMBEDDING_MODEL)
    qdrant_manager.init_collection(st.session_state.encoder)

if "messages" not in st.session_state:
    try:
        st.session_state.messages = message_manager.load_messages(qdrant_manager)
    except Exception as e:
        logging.warning(f"Не удалось загрузить сообщения: {e}")
        st.session_state.messages = []

if "ollama_status" not in st.session_state:
    st.session_state.ollama_status, st.session_state.ollama_message = ollama_client.check_api(config.OLLAMA_URL)

if "context_cache" not in st.session_state:
    st.session_state.context_cache = OrderedDict()
    st.session_state.context_cache_max_size = config.MAX_CACHE_SIZE

if "stats" not in st.session_state:
    st.session_state.stats = {
        "messages_sent": 0,
        "files_processed": 0,
        "chunks_saved": 0
    }

# ==================== ASYNC THREAD HELPER ====================
def async_stream_thread(q: Queue, *args):
    async def inner_gen():
        try:
            async for token in ollama_client.ask_llm_async_stream(*args):
                q.put(token)
        except Exception as e:
            logging.error(f"Ошибка в inner_gen: {e}")
            q.put(f"⚠️ Ошибка стрима: {e}")
        finally:
            q.put(None)

    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    try:
        loop.run_until_complete(inner_gen())
    finally:
        loop.close()

# ==================== UI COMPONENTS ====================
def render_sidebar():
    with st.sidebar:
        st.header("⚙️ Настройки")
        base_url = st.text_input("Ollama URL", config.OLLAMA_URL)
        model = st.text_input("Модель", config.DEFAULT_LLM_MODEL)
        max_context = st.slider("Контекст для ответа", 2, 20, 6)
        embedder = st.selectbox("Модель эмбеддингов", ["all-MiniLM-L6-v2", "all-mpnet-base-v2"])
        max_tokens = st.slider("Максимальная длина ответа", 50, 1000, 200)
        temperature = st.slider("Температура", 0.0, 1.5, 0.3)
        top_p = st.slider("Top-p", 0.1, 1.0, 0.8)
        top_k = st.slider("Top-k", 1, 100, 40)

        if st.button("🔄 Перезагрузить энкодер"):
            st.session_state.encoder = encoder_manager.load_encoder(embedder)
            qdrant_manager.init_collection(st.session_state.encoder)
            st.success(f"Энкодер {embedder} загружен")

        # Определяем поддерживаемые форматы динамически
        supported_formats = ["txt", "pdf", "docx", "csv"]
        if PPTX_AVAILABLE:
            supported_formats.append("pptx")
        if XLSX_AVAILABLE:
            supported_formats.append("xlsx")
            
        uploaded_files = st.file_uploader(
            "📂 Загрузите файлы", 
            type=supported_formats, 
            accept_multiple_files=True
        )
        
        st.markdown("---")
        if st.button("💾 Экспорт чата"):
            if st.session_state.messages:
                chat_data = {
                    "exported_at": datetime.now().isoformat(),
                    "messages": st.session_state.messages,
                    "stats": st.session_state.stats
                }
                st.download_button(
                    label="📥 Скачать историю",
                    data=json.dumps(chat_data, ensure_ascii=False, indent=2),
                    file_name=f"chat_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                    mime="application/json"
                )
            else:
                st.info("Нет сообщений для экспорта")

        uploaded_chat = st.file_uploader("📤 Импорт чата", type=["json"])
        if uploaded_chat:
            try:
                chat_data = json.load(uploaded_chat)
                st.session_state.messages = chat_data.get("messages", [])
                if "stats" in chat_data:
                    st.session_state.stats = chat_data["stats"]
                st.success("История чата импортирована")
                st.rerun()
            except Exception as e:
                st.error(f"Ошибка импорта: {e}")
        
        if st.button("🗑 Очистить чат"):
            st.session_state.messages = []
            try:
                qdrant_manager.clear_collection(st.session_state.encoder)
                st.session_state.context_cache.clear()
                encoder_manager.cache.clear()
                st.session_state.stats = {
                    "messages_sent": 0,
                    "files_processed": 0,
                    "chunks_saved": 0
                }
                st.success("История чата очищена")
                logging.info("История чата очищена")
            except Exception as e:
                st.error(f"Ошибка при очистке Qdrant: {e}")
                logging.error(f"Ошибка при очистке Qdrant: {e}")
            st.rerun()

        st.markdown("---")
        if st.button("📂 Показать загруженные файлы"):
            try:
                file_messages = []
                all_messages = qdrant_manager.load_messages(limit=100)
                for msg in all_messages:
                    if msg.get('role') == 'file':
                        file_messages.append(msg)
                
                if file_messages:
                    st.write(f"Загружено файлов: {len(set(msg.get('time', '') for msg in file_messages))}")
                    files_by_time = {}
                    for msg in file_messages:
                        time_key = msg.get('time', 'unknown')
                        if time_key not in files_by_time:
                            files_by_time[time_key] = []
                        files_by_time[time_key].append(msg)
                    
                    for time_key, chunks in files_by_time.items():
                        with st.expander(f"Файл от {time_key} ({len(chunks)} чанков)"):
                            if chunks:
                                st.text_area("Содержимое:", chunks[0].get('content', '')[:500] + "...", height=100, key=f"file_{time_key}")
                else:
                    st.info("Нет загруженных файлов")
            except Exception as e:
                st.error(f"Ошибка при получении файлов: {e}")

        st.markdown("---")
        st.subheader("📊 Статистика")
        st.write(f"Сообщений: {st.session_state.stats.get('messages_sent', 0)}")
        st.write(f"Файлов: {st.session_state.stats.get('files_processed', 0)}")
        st.write(f"Чанков: {st.session_state.stats.get('chunks_saved', 0)}")
        # Расширенная статистика кэша
        cache_stats = encoder_manager.cache.get_stats()
        st.write(f"Кэш эмбеддингов: {cache_stats['size']}/{cache_stats['max_size']} ({cache_stats['usage_percent']:.1f}%)")

    return base_url, model, max_context, max_tokens, temperature, top_p, top_k, uploaded_files

# ==================== MAIN ====================
def main():
    st.set_page_config(page_title="LLM Chat", layout="wide")
    st.title("💬 LLM Чат")

    if config.DEBUG_MODE:
        st.sidebar.info("DEBUG режим включен")

    base_url, model, max_context, max_tokens, temperature, top_p, top_k, uploaded_files = render_sidebar()

    if st.session_state.ollama_status:
        st.success(f"🟢 {st.session_state.ollama_message}")
    else:
        st.error(f"🔴 {st.session_state.ollama_message}")

    if uploaded_files:
        for file in uploaded_files:
            result = file_processor.process_file(file, qdrant_manager, encoder_manager, config.MAX_FILE_SIZE_MB)
            st.sidebar.success(result)

    for msg in st.session_state.messages:
        with st.chat_message(msg.get("role", "user")):
            st.write(msg.get("content", ""))
            st.caption(f"_{msg.get('time', '')}_")

    if st.session_state.ollama_status and (prompt := st.chat_input("Введите сообщение...")):
        if not validator.validate_prompt(prompt):
            st.rerun()
        
        ts = datetime.now().strftime("%H:%M:%S")
        qdrant_manager.save_message("user", prompt, ts, encoder_manager)
        st.session_state.messages.append({"role": "user", "content": prompt, "time": ts})
        
        st.session_state.stats["messages_sent"] = st.session_state.stats.get("messages_sent", 0) + 1
        
        with st.chat_message("user"):
            st.write(prompt)
            st.caption(f"_{ts}_")

        with st.chat_message("assistant"):
            # Индикатор прогресса для генерации
            progress_placeholder = st.empty()
            progress_placeholder.info("⏳ Подготовка ответа...")
            
            with st.spinner("Генерация ответа..."):
                context = qdrant_manager.search_context(prompt, encoder_manager, max_context)
                progress_placeholder.empty()  # Убираем индикатор
                
                if config.DEBUG_MODE:
                    with st.expander("🔍 Реальный контекст (отладка)", expanded=False):
                        st.write(f"Найдено сообщений: {len(context)}")
                        chat_count = sum(1 for msg in context if msg.get('role') in ['user', 'assistant'])
                        file_count = sum(1 for msg in context if msg.get('role') == 'file')
                        st.write(f"Чат: {chat_count}, Файлы: {file_count}")
                        
                        for i, msg in enumerate(context):
                            role = msg.get('role', 'unknown')
                            role_emoji = "👤" if role == "user" else "🤖" if role == "assistant" else "📄" if role == "file" else "❓"
                            st.write(f"**{i+1}. {role_emoji} {role.capitalize()}** ({msg.get('time', 'N/A')}):")
                            content = msg.get('content', '')
                            st.code(content[:300] + ('...' if len(content) > 300 else ''))
                
                stop_button = st.empty()
                stop_generation = False
                
                q = Queue()
                thread = threading.Thread(
                    target=async_stream_thread,
                    args=(q, prompt, context, model, base_url, max_tokens, temperature, top_p, top_k),
                    daemon=True
                )
                thread.start()

                full_response = ''
                placeholder = st.empty()
                
                stop_key = f"stop_{ts.replace(':', '_')}"
                if stop_button.button("⏹️ Остановить генерацию", key=stop_key):
                    stop_generation = True
                    stop_button.empty()
                
                # Таймер для таймаута
                start_time = datetime.now()
                
                while not stop_generation:
                    try:
                        chunk = q.get(timeout=0.5)  # Увеличенный таймаут
                        if chunk is None:
                            break
                        full_response += chunk
                        placeholder.markdown(full_response)
                        
                        # Проверка таймаута
                        if (datetime.now() - start_time).seconds > Config.STREAM_TIMEOUT:
                            stop_generation = True
                            full_response += "\n\n⚠️ Превышено время ожидания"
                            break
                            
                    except:
                        # Таймаут - проверяем нужно ли остановиться
                        if stop_generation or (datetime.now() - start_time).seconds > Config.STREAM_TIMEOUT:
                            break
                        continue
                
                # Улучшенное ожидание завершения потока
                try:
                    thread.join(timeout=5)  # Увеличенный таймаут
                except:
                    pass
                stop_button.empty()

        ts_reply = datetime.now().strftime("%H:%M:%S")
        if stop_generation and not full_response.strip():
            final_response = "Генерация была остановлена пользователем."
        elif not full_response.strip():
            final_response = "Извините, не удалось сгенерировать ответ."
        else:
            final_response = full_response.strip()

        qdrant_manager.save_message("assistant", final_response, ts_reply, encoder_manager)
        st.session_state.messages.append({"role": "assistant", "content": final_response, "time": ts_reply})
        st.rerun()

if __name__ == "__main__":
    main()