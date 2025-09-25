import io
import csv
import uuid
import logging
import asyncio
import streamlit as st
from typing import List
from datetime import datetime
from PyPDF2 import PdfReader
from docx import Document

# Добавляем fallback для PDF и других библиотек
try:
    import pdfplumber
    PDF_PLUMBER_AVAILABLE = True
except ImportError:
    PDF_PLUMBER_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

from qdrant_client.http import models as qmodels
from app.config import Config

class FileProcessor:
    def __init__(self):
        pass

    def chunk_text(self, text: str, chunk_size: int = Config.CHUNK_SIZE, 
                   overlap: int = Config.CHUNK_OVERLAP) -> List[str]:
        if not text:
            return []
        chunks = []
        start = 0
        text_len = len(text)
        while start < text_len:
            end = min(start + chunk_size, text_len)
            chunks.append(text[start:end])
            start += chunk_size - overlap
            if start >= text_len:
                break
        return chunks

    def process_file(self, file, qdrant_manager, encoder_manager, max_size_mb: int = Config.MAX_FILE_SIZE_MB) -> str:
        file_ext = file.name.lower().split('.')[-1] if '.' in file.name else ""
        allowed_extensions = {"txt", "pdf", "docx", "csv"}
        
        # Добавляем поддержку новых форматов только если библиотеки доступны
        if PPTX_AVAILABLE:
            allowed_extensions.add("pptx")
        if XLSX_AVAILABLE:
            allowed_extensions.add("xlsx")
        
        if not file_ext:
            return f"⚠️ Файл без расширения не поддерживается"
        
        if file_ext not in allowed_extensions:
            available_formats = ", ".join(sorted(allowed_extensions))
            return f"⚠️ Формат .{file_ext} не поддерживается. Доступные: {available_formats}"
        
        if file.size > max_size_mb * 1024 * 1024:
            return f"⚠️ Файл '{file.name}' слишком большой (макс. {max_size_mb} МБ)"

        try:
            file.seek(0)
            raw_bytes = file.read()
            file.seek(0)
        except Exception as e:
            return f"⚠️ Не удалось прочитать файл {file.name}: {e}"

        content = ""
        try:
            if file_ext == "txt":
                content = raw_bytes.decode("utf-8", errors="ignore")
            elif file_ext == "pdf":
                content = self._extract_pdf_content(raw_bytes)
            elif file_ext == "docx":
                content = self._extract_docx_content(raw_bytes)
            elif file_ext == "csv":
                content = self._extract_csv_content(raw_bytes)
            elif file_ext == "pptx":
                if PPTX_AVAILABLE:
                    content = self._extract_pptx_content(raw_bytes)
                else:
                    content = "⚠️ Обработка PPTX не поддерживается (отсутствует библиотека python-pptx)"
            elif file_ext == "xlsx":
                if XLSX_AVAILABLE:
                    content = self._extract_xlsx_content(raw_bytes)
                else:
                    content = "⚠️ Обработка XLSX не поддерживается (отсутствует библиотека openpyxl)"
            else:
                content = f"⚠️ Обработка формата .{file_ext} не поддерживается"
        except Exception as e:
            content = f"⚠️ Ошибка при обработке файла: {e}"
            logging.error(f"Ошибка обработки файла {file.name}: {e}")

        chunks = self.chunk_text(content)
        if chunks:
            self._save_chunks_as_messages(chunks, datetime.now().strftime("%H:%M:%S"), qdrant_manager, encoder_manager)
            if "stats" in st.session_state:
                st.session_state.stats["files_processed"] = st.session_state.stats.get("files_processed", 0) + 1
                st.session_state.stats["chunks_saved"] = st.session_state.stats.get("chunks_saved", 0) + len(chunks)
            return f"✅ Файл '{file.name}' загружен ({len(chunks)} чанков)"
        else:
            return f"⚠️ Файл '{file.name}' пустой или нераспознанный"

    def _extract_pdf_content(self, raw_bytes) -> str:
        # Пробуем PyPDF2
        try:
            reader = PdfReader(io.BytesIO(raw_bytes))
            total_pages = len(reader.pages)

            if total_pages == 0:
                return ""
            
            progress_bar = st.progress(0) if total_pages > 2 else None
            
            content_parts = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    content_parts.append(text + "\n")
                else:
                    # Fallback: пробуем pdfplumber если страница пустая
                    if PDF_PLUMBER_AVAILABLE:
                        try:
                            with pdfplumber.open(io.BytesIO(raw_bytes)) as pdf:
                                if i < len(pdf.pages):
                                    page_obj = pdf.pages[i]
                                    text_alt = page_obj.extract_text()
                                    if text_alt:
                                        content_parts.append(text_alt + "\n")
                        except Exception:
                            pass
                
                if progress_bar and i % max(1, total_pages // 20) == 0:
                    progress_bar.progress((i + 1) / total_pages)
            
            if progress_bar:
                progress_bar.empty()
            
            return "".join(content_parts)
        except Exception as e:
            logging.warning(f"PyPDF2 ошибка: {e}")
            # Fallback на pdfplumber если доступен
            if PDF_PLUMBER_AVAILABLE:
                try:
                    content_parts = []
                    with pdfplumber.open(io.BytesIO(raw_bytes)) as pdf:
                        progress_bar = st.progress(0) if len(pdf.pages) > 2 else None
                        for i, page in enumerate(pdf.pages):
                            text = page.extract_text()
                            if text:
                                content_parts.append(text + "\n")
                            if progress_bar and i % max(1, len(pdf.pages) // 20) == 0:
                                progress_bar.progress((i + 1) / len(pdf.pages))
                        if progress_bar:
                            progress_bar.empty()
                    return "".join(content_parts)
                except Exception as e2:
                    logging.error(f"pdfplumber ошибка: {e2}")
            return f"⚠️ Ошибка извлечения PDF: {e}"

    def _extract_docx_content(self, raw_bytes) -> str:
        doc = Document(io.BytesIO(raw_bytes))
        content_parts = [para.text + "\n" for para in doc.paragraphs]
        return "".join(content_parts)

    def _extract_csv_content(self, raw_bytes) -> str:
        decoded = io.StringIO(raw_bytes.decode("utf-8", errors="ignore"))
        reader = csv.reader(decoded)
        return "\n".join([", ".join(row) for row in reader])

    def _extract_pptx_content(self, raw_bytes) -> str:
        if not PPTX_AVAILABLE:
            return "⚠️ Библиотека python-pptx не установлена"
            
        try:
            prs = Presentation(io.BytesIO(raw_bytes))
            content_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                content_parts.append(f"Slide {slide_num}:\n")
                for shape in slide.shapes:
                    # Извлекаем текст из text_frame, если он есть
                    if getattr(shape, "has_text_frame", False):
                        text_frame = getattr(shape, "text_frame", None)
                        if text_frame:
                            for paragraph in getattr(text_frame, "paragraphs", []):
                                text = getattr(paragraph, "text", "").strip()
                                if text:
                                    content_parts.append(text + "\n")
                    else:
                        text = getattr(shape, "text", "").strip()
                        if text:
                            content_parts.append(text + "\n")

            return "\n".join(content_parts)
        except Exception as e:
            return f"⚠️ Ошибка извлечения PPTX: {e}"

    def _extract_xlsx_content(self, raw_bytes) -> str:
        if not XLSX_AVAILABLE:
            return "⚠️ Библиотека openpyxl не установлена"
            
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True)
            content_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_parts.append(f"Sheet: {sheet_name}\n")
                for row in sheet.iter_rows(values_only=True):
                    content_parts.append(", ".join(str(cell) if cell is not None else "" for cell in row) + "\n")
            return "".join(content_parts)
        except Exception as e:
            return f"⚠️ Ошибка извлечения XLSX: {e}"

    async def _save_chunks_as_messages_async(self, chunks: list, timestamp: str, 
                                           qdrant_manager, encoder_manager, 
                                           role: str = "file", progress_bar=None):
        try:
            if not chunks:
                return
                
            vectors = await encoder_manager.async_encode(chunks, batch_size=min(32, len(chunks)))
            
            points = []
            for chunk, vec in zip(chunks, vectors):
                vector_list = vec if isinstance(vec, list) else list(vec)
                point = qmodels.PointStruct(
                    id=str(uuid.uuid4()),
                    vector=vector_list,
                    payload={"role": role, "content": chunk, "time": timestamp}
                )
                points.append(point)
            
            batch_size = 64
            total_points = len(points)
            
            for i in range(0, total_points, batch_size):
                batch = points[i:i + batch_size]
                qdrant_manager.client.upsert(
                    collection_name=qdrant_manager.collection_name, 
                    points=batch
                )
                
                if progress_bar and total_points > batch_size:
                    progress_bar.progress(min(1.0, (i + batch_size) / total_points))
            
            logging.info(f"Сохранено {len(points)} чанков текста (асинхронно)")
            
        except Exception as e:
            logging.error(f"Ошибка асинхронного сохранения чанков: {e}")
            raise

    def _save_chunks_as_messages(self, chunks: list, timestamp: str, 
                               qdrant_manager, encoder_manager, role: str = "file"):
        try:
            try:
                loop = asyncio.get_event_loop()
            except RuntimeError:
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)
            
            progress = st.progress(0) if len(chunks) > 32 else None
            
            # Используем asyncio.wait_for для таймаута
            try:
                future = asyncio.wait_for(
                    self._save_chunks_as_messages_async(chunks, timestamp, qdrant_manager, encoder_manager, role, progress),
                    timeout=Config.STREAM_TIMEOUT
                )
                loop.run_until_complete(future)
            except asyncio.TimeoutError:
                st.sidebar.warning("⚠️ Таймаут при сохранении чанков")
                logging.error("Таймаут при сохранении чанков")
            
            if progress:
                progress.empty()
                
        except Exception as e:
            logging.error(f"Ошибка сохранения чанков: {e}")
            st.sidebar.warning(f"Ошибка сохранения чанков: {e}")